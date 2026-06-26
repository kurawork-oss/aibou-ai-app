import urllib.parse
import json
import re
import io
import base64
from gtts import gTTS
from streamlit_mic_recorder import speech_to_text
import google.generativeai as genai
import os
import streamlit as st

# 🚨 パワポ生成用のライブラリ
try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    st.error("⚠️ `python-pptx` ライブラリがインストールされていません。requirements.txt を確認してください。")

# 💎 UIデザイン用CSS
st.markdown("""
    <style>
    /* 1. 全体をダーク＆サイバーパンクな雰囲気に */
    [data-testid="stAppViewContainer"] {
        background-color: #030b14 !important;
        background-image: radial-gradient(circle at 50% 120%, rgba(150, 200, 255, 0.15), transparent) !important;
        color: #e2e8f0 !important;
    }

    /* 2. ⬡と❖のシンボルを光らせる */
    .saas-title {
        color: #ffffff !important;
        font-weight: 900;
        letter-spacing: 6px;
        margin-bottom: 5px;
        text-shadow: 0 0 10px rgba(197, 198, 199, 0.6), 0 0 20px rgba(150, 200, 255, 0.4) !important;
        text-align: center;
    }
    .central-logo {
        text-align: center;
        font-size: 70px;
        color: #c5c6c7 !important;
        text-shadow: 0 0 15px rgba(197, 198, 199, 0.8), 0 0 30px rgba(150, 200, 255, 0.5) !important;
        margin: 30px 0 10px 0;
    }
    .central-logo-sub {
        text-align: center;
        color: #a0aec0;
        font-size: 12px;
        letter-spacing: 4px;
        margin-top: -10px;
        margin-bottom: 40px;
    }

    /* 3. ボタンのダークサイバー仕様 */
    div.stButton > button {
        background-color: rgba(15, 23, 42, 0.8) !important; 
        border: 1px solid rgba(255, 255, 255, 0.2) !important; 
        border-radius: 12px !important;
        transition: all 0.3s ease !important;
        padding: 10px !important;
        height: 60px !important;
    }
    div.stButton > button p {
        color: #ffffff !important;
        font-weight: 800 !important;
        letter-spacing: 2px !important;
        margin: 0 !important;
    }

    /* 4. ホバー時の白いネオン発光エフェクト */
    div.stButton > button:hover {
        background-color: rgba(30, 41, 59, 0.9) !important;
        border-color: #ffffff !important;
        box-shadow: 0 0 20px rgba(255, 255, 255, 0.6), inset 0 0 10px rgba(255, 255, 255, 0.2) !important;
        transform: translateY(-3px) !important;
    }
    div.stButton > button:hover p {
        text-shadow: 0 0 10px #ffffff, 0 0 20px #ffffff !important;
    }
    
    /* 5. ホバー時に下部に表示される説明エリア */
    .desc-display-area {
        position: relative;
        height: 150px;
        border: 1px solid rgba(255, 255, 255, 0.1);
        border-radius: 12px;
        margin-top: 50px;
        display: flex;
        justify-content: center;
        align-items: center;
        background: rgba(10, 20, 40, 0.4);
        transition: all 0.3s ease;
        overflow: hidden;
    }
    
    .desc-text {
        position: absolute;
        width: 100%;
        text-align: center;
        opacity: 0;
        transition: opacity 0.3s ease, transform 0.3s ease;
        transform: translateY(10px);
        color: #ffffff;
        font-size: 14px;
        font-weight: bold;
        letter-spacing: 1px;
        line-height: 1.5;
        text-shadow: 0 0 10px rgba(255, 255, 255, 0.3);
        pointer-events: none;
    }
    
    .default-desc {
        opacity: 1;
        transform: translateY(0);
        color: #718096;
        letter-spacing: 2px;
    }

    /* 魔法のCSS：Streamlitのカラムの順番でホバーを「確実」に検知する */
    .stApp:has([data-testid="column"]:nth-of-type(1) button:hover) .app-desc,
    .stApp:has([data-testid="stColumn"]:nth-of-type(1) button:hover) .app-desc { opacity: 1 !important; transform: translateY(0) !important; }
    
    .stApp:has([data-testid="column"]:nth-of-type(2) button:hover) .img-desc,
    .stApp:has([data-testid="stColumn"]:nth-of-type(2) button:hover) .img-desc { opacity: 1 !important; transform: translateY(0) !important; }
    
    .stApp:has([data-testid="column"]:nth-of-type(3) button:hover) .vid-desc,
    .stApp:has([data-testid="stColumn"]:nth-of-type(3) button:hover) .vid-desc { opacity: 1 !important; transform: translateY(0) !important; }
    
    .stApp:has([data-testid="column"]:nth-of-type(4) button:hover) .slide-desc,
    .stApp:has([data-testid="stColumn"]:nth-of-type(4) button:hover) .slide-desc { opacity: 1 !important; transform: translateY(0) !important; }

    .stApp:has(button:hover) .default-desc { opacity: 0 !important; transform: translateY(-10px) !important; }
    .stApp:has([data-testid="stHorizontalBlock"] button:hover) .desc-display-area {
        border-color: #ffffff !important;
        box-shadow: 0 0 20px rgba(197, 198, 199, 0.4), inset 0 0 10px rgba(197, 198, 199, 0.2) !important;
        background: rgba(15, 30, 50, 0.8) !important;
    }

    /* 下部の光るプラットフォーム */
    .hologram-platform {
        position: fixed;
        bottom: -80px;
        left: 50%;
        transform: translateX(-50%);
        width: 600px;
        height: 150px;
        background: radial-gradient(ellipse at 50% 50%, rgba(197, 198, 199, 0.15), transparent 70%);
        border-radius: 50%;
        box-shadow: 0 0 50px rgba(197, 198, 199, 0.2);
        pointer-events: none;
        z-index: -1;
    }
    
    /* 🌟 【絶対解決版】プロジェクトカード（強烈なセレクタ＋ハッキリ見える色） */
    div[data-testid="stVerticalBlockBorderWrapper"] {
        background-color: #1e293b !important; /* ← 絶対に漆黒から浮き出るスレートグレー！ */
        border: 2px solid rgba(197, 198, 199, 0.5) !important; /* ← 確実に見える2pxのシアン枠！ */
        border-radius: 16px !important;
        box-shadow: 0 10px 25px rgba(0, 0, 0, 0.8) !important; /* 濃い影で浮かせる */
        transition: all 0.3s ease !important;
    }

    /* 🌟 プロジェクトカードにマウスを乗せた時のホバーエフェクト */
    div[data-testid="stVerticalBlockBorderWrapper"]:hover {
        background-color: #334155 !important; /* ホバーでさらに一段階明るく！ */
        border-color: #c5c6c7 !important; /* ネオンシアン発光 */
        box-shadow: 0 15px 35px rgba(0, 0, 0, 0.9), 0 0 20px rgba(197, 198, 199, 0.4) !important;
        transform: translateY(-5px) !important;
    }
    </style>
""", unsafe_allow_html=True)

# 🚨 初期化
if "auto_fix_prompt" not in st.session_state: st.session_state.auto_fix_prompt = ""
if "forge_workspaces" not in st.session_state: st.session_state.forge_workspaces = {}
if "current_forge_ws" not in st.session_state: st.session_state.current_forge_ws = None 
if "ai_voice_base64" not in st.session_state: st.session_state.ai_voice_base64 = None
if "just_generated_audio" not in st.session_state: st.session_state.just_generated_audio = False
if "selected_forge_mode" not in st.session_state: st.session_state.selected_forge_mode = None

# ==========================================
# 🚪 ステージ1：ホログラムカード選択画面
# ==========================================
if st.session_state.current_forge_ws is None and st.session_state.selected_forge_mode is None:
    st.markdown('<div class="central-logo">⬡</div>', unsafe_allow_html=True)
    st.markdown("<h2 class='saas-title'>❖ FORGE STUDIO ❖</h2>", unsafe_allow_html=True)
    st.markdown("<p class='central-logo-sub'>SELECT SYSTEM ENGINE</p>", unsafe_allow_html=True)
    room_help("Forge Lab")
    
    st.markdown('<div class="hologram-platform"></div>', unsafe_allow_html=True)

    r1 = st.columns(3, gap="large")
    with r1[0]:
        if st.button("APP STUDIO", use_container_width=True):
            st.session_state.selected_forge_mode = "APP"; st.rerun()
    with r1[1]:
        if st.button("IMAGE GENERATOR", use_container_width=True):
            st.session_state.selected_forge_mode = "IMAGE"; st.rerun()
    with r1[2]:
        if st.button("VIDEO PRODUCTION", use_container_width=True):
            st.session_state.selected_forge_mode = "VIDEO"; st.rerun()
    r2 = st.columns(3, gap="large")
    with r2[0]:
        if st.button("SLIDE DECK", use_container_width=True):
            st.session_state.selected_forge_mode = "SLIDE"; st.rerun()
    with r2[1]:
        if st.button("SHEET BUILDER", use_container_width=True):
            st.session_state.selected_forge_mode = "SHEET"; st.rerun()
    with r2[2]:
        if st.button("DOC WRITER", use_container_width=True):
            st.session_state.selected_forge_mode = "DOC"; st.rerun()

    st.markdown("""
        <div class="desc-display-area">
            <div class="desc-text default-desc">SELECT AN ENGINE / エンジンを選択（アプリ・画像・動画・スライド・表計算・文書）</div>
        </div>
    """, unsafe_allow_html=True)
    st.caption("🤖 APP=アプリ生成・実行／🎨 IMAGE=画像生成／🎬 VIDEO=絵コンテ＋MP4／"
               "📊 SLIDE=.pptx／📈 SHEET=表データ(CSV/Excel)／📝 DOC=文書(Markdown/Word)")

# ==========================================
# 🚪 ステージ2：モード別プロジェクト管理画面
# ==========================================
elif st.session_state.current_forge_ws is None and st.session_state.selected_forge_mode is not None:
    mode = st.session_state.selected_forge_mode
    
    col_back, col_title = st.columns([3, 7])
    with col_back:
        if st.button("← BACK TO ENGINES", use_container_width=True):
            st.session_state.selected_forge_mode = None
            st.rerun()
    with col_title:
        st.markdown(f"<h3 style='color:#ffffff; font-weight:800; margin-top:-5px;'>[ {mode} ENGINE ]</h3>", unsafe_allow_html=True)
    
    st.markdown("---")
    
    col_create, col_list = st.columns([4, 6], gap="large")
    
    with col_create:
        st.markdown("<p style='font-weight:bold; color:#a0aec0; font-size:12px;'>[ INITIALIZE NEW PROJECT ]</p>", unsafe_allow_html=True)
        with st.container(border=True):
            new_ws_name = st.text_input("PROJECT NAME", label_visibility="collapsed", placeholder="Enter project name...")
            st.markdown("<br>", unsafe_allow_html=True)
            if st.button("CREATE WORKSPACE", type="primary", use_container_width=True):
                if new_ws_name and new_ws_name not in st.session_state.forge_workspaces:
                    st.session_state.forge_workspaces[new_ws_name] = {
                        "type": mode, "chat": [], "code": "", "media": "", "retries": 0
                    }
                    st.session_state.current_forge_ws = new_ws_name
                    st.rerun()
                elif not new_ws_name:
                    st.error("Please enter a project name.")

    with col_list:
        st.markdown("<p style='font-weight:bold; color:#a0aec0; font-size:12px;'>[ ACTIVE WORKSPACES ]</p>", unsafe_allow_html=True)
        mode_workspaces = {k: v for k, v in st.session_state.forge_workspaces.items() if v.get("type") == mode}
        
        if not mode_workspaces:
            st.info("No active projects found for this engine.")
        else:
            cols = st.columns(2)
            for idx, (ws_name, ws_data) in enumerate(mode_workspaces.items()):
                with cols[idx % 2]:
                    with st.container(border=True):
                        st.markdown(f"<span style='border:1px solid #c5c6c7; color:#c5c6c7; padding:2px 8px; border-radius:4px; font-size:10px;'>{ws_data['type']}</span>", unsafe_allow_html=True)
                        st.markdown(f"<h4 style='color:#ffffff; font-weight:800; margin-top:10px;'>{ws_name}</h4>", unsafe_allow_html=True)
                        st.markdown(f"<p style='font-size: 11px; color: #718096;'>ONLINE | Logs: {len(ws_data['chat'])}</p>", unsafe_allow_html=True)
                        
                        c_btn1, c_btn2 = st.columns([7, 3])
                        with c_btn1:
                            if st.button("ENTER", key=f"open_{ws_name}", use_container_width=True):
                                st.session_state.current_forge_ws = ws_name
                                st.rerun()
                        with c_btn2:
                            if st.button("DEL", key=f"del_{ws_name}", use_container_width=True):
                                del st.session_state.forge_workspaces[ws_name]
                                st.rerun()

# ==========================================
# 🖥️ ステージ3：ワークスペース内部画面
# ==========================================
else:
    ws_name = st.session_state.current_forge_ws
    ws_data = st.session_state.forge_workspaces[ws_name]
    ws_type = ws_data.get("type", "APP")
    
    if "retries" not in ws_data: ws_data["retries"] = 0
    if "media" not in ws_data: ws_data["media"] = ""
    
    col_back, col_title = st.columns([3, 7])
    with col_back:
        if st.button("← BACK TO PROJECT LIST", use_container_width=True):
            st.session_state.current_forge_ws = None
            st.rerun()
    with col_title:
        st.markdown(f"<h3 style='color:#ffffff; font-weight:800; margin-top:-5px;'>[ {ws_name} ] <span style='font-size:14px; color:#718096; font-weight:normal;'>| ENGINE: {ws_type}</span></h3>", unsafe_allow_html=True)

    st.markdown("---")

    with st.sidebar:
        st.markdown(f"<div style='text-align:center; font-weight:800; color:#c5c6c7; margin-bottom:10px;'>[ {ws_name} ]</div>", unsafe_allow_html=True)
        with st.form("forge_sidebar_form", clear_on_submit=True):
            placeholder_text = "Type your prompt here..."
            if ws_type == "APP": placeholder_text = "例：シンプルな計算機アプリ"
            elif ws_type == "IMAGE": placeholder_text = "例：サイバーパンクな都市"
            elif ws_type == "VIDEO": placeholder_text = "例：コーヒーが弾ける動画"
            elif ws_type == "SLIDE": placeholder_text = "例：AIの未来について5枚"
            elif ws_type == "SHEET": placeholder_text = "例：2024年の月別売上表（商品別）"
            elif ws_type == "DOC": placeholder_text = "例：新サービスの企画書を作成"
            
            forge_prompt = st.text_area("PROMPT", placeholder=placeholder_text, height=150, label_visibility="collapsed")
            submitted = st.form_submit_button("EXECUTE", use_container_width=True, type="primary")
        
        # 🎙️ 音声入力は安定性のため既定オフ（旧コンポーネント）。HUBで有効化したときのみ。
        spoken_text = None
        if st.session_state.get("mic_enabled"):
            try:
                spoken_text = speech_to_text(language='ja', start_prompt="🎙️ VOICE COMMAND", stop_prompt="🛑 SEND", use_container_width=True, just_once=True, key='Forge_STT')
            except Exception:
                spoken_text = None

    col_log, col_preview = st.columns([3, 7])
    
    with col_log:
        st.markdown("<p style='font-weight:bold; color:#a0aec0; font-size:12px;'>[ SYSTEM CONSOLE ]</p>", unsafe_allow_html=True)
        _fresh = bool(st.session_state.get("just_generated_audio"))
        st.session_state.just_generated_audio = False
        render_core(200)
        if st.session_state.get("ai_voice_base64") and _fresh:
            try:
                st.audio(base64.b64decode(st.session_state.ai_voice_base64), format="audio/mp3", autoplay=True)
            except Exception:
                pass

        with st.container(height=400, border=False):
            if not ws_data["chat"]:
                st.info("System Ready. Waiting for input...")
            for m in ws_data["chat"]:
                with st.chat_message(m["role"], avatar="👤" if m["role"]=="user" else "🤖"):
                    st.markdown(m["content"])

    with col_preview:
        st.markdown("<p style='font-weight:bold; color:#a0aec0; font-size:12px;'>[ CANVAS PREVIEW ]</p>", unsafe_allow_html=True)
                
        if st.session_state.auto_fix_prompt:
            st.warning("⚙️ RECOVERY PROTOCOL INITIATED: AI is fixing the code...")
            
        elif ws_type == "APP":
            if ws_data["code"]:
                st.download_button(label="[ DOWNLOAD .py ]", data=ws_data["code"], file_name=f"{ws_name.replace(' ', '_')}.py", mime="text/plain", use_container_width=True)
                with st.container(border=True):
                    try:
                        exec(ws_data["code"], globals())
                        ws_data["retries"] = 0
                    except Exception as e:
                        st.error(f"RUNTIME ERROR:\n{e}")
                        if ws_data["retries"] < 3:
                            ws_data["retries"] += 1
                            st.session_state.auto_fix_prompt = f"実行時に以下のエラーが発生しました。修正して！\n\n【エラー内容】\n{e}"
                            st.rerun()
                        else:
                            st.error("❌ 自己修復が上限に達しました。")
                with st.expander("📝 MANUAL OVERRIDE (SOURCE CODE)"):
                    edited_code = st.text_area("Python Code", value=ws_data["code"], height=300)
                    if st.button("UPDATE CODE"):
                        ws_data["code"] = edited_code
                        st.rerun()
            else:
                st.info("Canvas is empty.")
                
        elif ws_type == "IMAGE":
            if ws_data["media"]:
                image_url = f"https://image.pollinations.ai/prompt/{urllib.parse.quote(ws_data['media'])}?width=800&height=450&nologo=true"
                st.image(image_url, use_container_width=True)
                
                with st.expander("⚙️ PROMPT TUNING (ENGLISH)"):
                    edited_media = st.text_area("Prompt", value=ws_data["media"], height=100)
                    if st.button("UPDATE IMAGE"):
                        ws_data["media"] = edited_media
                        st.rerun()
            else:
                st.info("Canvas is empty.")
                
        elif ws_type == "VIDEO":
            if ws_data["code"]:
                st.success("✅ 絵コンテ／動画プロンプトが完成。画像＋ナレーションでMP4を生成できます。")

                if st.button("🎬 動画(MP4)を生成", use_container_width=True, type="primary"):
                    import renderer
                    _txt = ws_data["code"]
                    # 英語の動画プロンプト（バッククォート内）を抽出
                    _ep = ""
                    _m = re.search(r'プロンプト[^\n]*\n+`+\s*([^`]+?)\s*`+', _txt, re.DOTALL)
                    if not _m:
                        _m = re.search(r'`([^`]{20,})`', _txt)
                    if _m:
                        _ep = _m.group(1).strip()
                    # シーン（Scene N: 説明）を抽出。無ければコンセプト全体を1シーンに。
                    _scenes = []
                    for _line in _txt.split("\n"):
                        _sm = re.search(r'Scene\s*\d+[^:：]*[:：]\s*(.+)', _line)
                        if _sm:
                            _narr = re.sub(r'[*`#>_]', '', _sm.group(1)).strip(" -")
                            if _narr:
                                _scenes.append({"narration": _narr, "visual": _ep})
                    if not _scenes:
                        _concept = re.sub(r'[*`#>_]', '', _txt)[:500].strip()
                        _scenes = [{"narration": _concept or "AI generated video", "visual": _ep}]

                    if not renderer.is_available():
                        st.error("⚠️ 動画合成にはFFmpegが必要です。packages.txt に ffmpeg を追加済みのため、再デプロイ後に利用できます。")
                    else:
                        with st.spinner("🎬 画像とナレーションからMP4を合成中…（数十秒かかります）"):
                            _vp = renderer.render_forge_video(_scenes, _ep)
                        if _vp:
                            ws_data["video_path"] = _vp
                            st.rerun()
                        else:
                            st.error("動画の生成に失敗しました（画像取得・ネットワーク・FFmpegをご確認ください）。")

                if ws_data.get("video_path") and os.path.exists(ws_data["video_path"]):
                    st.video(ws_data["video_path"])
                    with open(ws_data["video_path"], "rb") as _vf:
                        st.download_button("[ DOWNLOAD .mp4 ]", data=_vf.read(),
                                           file_name=f"{ws_name.replace(' ', '_')}.mp4",
                                           mime="video/mp4", use_container_width=True)
                
                with st.expander("📝 STORYBOARD & PROMPT", expanded=True):
                    edited_code = st.text_area("Markdown Editor", value=ws_data["code"], height=300)
                    if st.button("UPDATE SCRIPT"):
                        ws_data["code"] = edited_code
                        st.rerun()
                st.markdown(ws_data["code"])
            else:
                st.info("Canvas is empty.")
                
        elif ws_type == "SLIDE":
            if ws_data["code"]:
                st.success("✅ プレゼン構成が完了しました。下のボタンから PowerPoint をダウンロードできます！")
                try:
                    prs = Presentation()
                    # AIはMarkdown(--- 区切り)を生成する。まずJSONを試し、ダメならMarkdownを解析。
                    try:
                        slides_data = json.loads(ws_data["code"])
                    except Exception:
                        slides_data = None
                    if not isinstance(slides_data, list):
                        slides_data = []
                        for _block in re.split(r'\n-{3,}\n', ws_data["code"]):
                            _block = _block.strip()
                            if not _block:
                                continue
                            _title, _bullets = "", []
                            for _ln in _block.split("\n"):
                                _s = _ln.strip()
                                if not _s:
                                    continue
                                _h = re.match(r'^#{1,3}\s+(.*)', _s)
                                if _h and not _title:
                                    _title = _h.group(1).strip()
                                else:
                                    _bullets.append(re.sub(r'^[\-\*>#\s]+', '', _s))
                            slides_data.append({"title": _title or "Slide", "bullets": _bullets})

                    for slide_info in slides_data:
                        slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(slide_layout)
                        title = slide.shapes.title
                        content = slide.placeholders[1]
                        
                        title.text = slide_info.get("title", "No Title")
                        tf = content.text_frame
                        for bullet in slide_info.get("bullets", []):
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.font.size = Pt(24)
                            
                    pptx_io = io.BytesIO()
                    prs.save(pptx_io)
                    pptx_io.seek(0)
                    
                    st.download_button(
                        label="[ DOWNLOAD .pptx ]",
                        data=pptx_io,
                        file_name=f"{ws_name.replace(' ', '_')}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        type="primary",
                        use_container_width=True
                    )
                except Exception as e:
                    st.error(f"パワポの生成に失敗しました: {e}")
                    st.code(ws_data["code"])
                
                with st.expander("⚙️ SLIDE DATA (Markdown / JSON)"):
                    edited_code = st.text_area("Slide Editor", value=ws_data["code"], height=300)
                    if st.button("UPDATE SLIDES"):
                        ws_data["code"] = edited_code
                        st.rerun()
            else:
                st.info("Canvas is empty.")

        elif ws_type == "SHEET":
            if ws_data["code"]:
                st.success("✅ 表データが完成しました。CSV / Excel でダウンロードできます。")
                _safe = ws_name.replace(' ', '_')
                try:
                    import pandas as _pd
                    _df = _pd.read_csv(io.StringIO(ws_data["code"]))
                    st.dataframe(_df, use_container_width=True)
                    st.download_button("[ DOWNLOAD .csv ]", data=ws_data["code"].encode("utf-8-sig"),
                                       file_name=f"{_safe}.csv", mime="text/csv", use_container_width=True)
                    try:
                        _xio = io.BytesIO()
                        with _pd.ExcelWriter(_xio, engine="openpyxl") as _w:
                            _df.to_excel(_w, index=False, sheet_name="Sheet1")
                        st.download_button("[ DOWNLOAD .xlsx ]", data=_xio.getvalue(),
                                           file_name=f"{_safe}.xlsx",
                                           mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                                           type="primary", use_container_width=True)
                    except Exception:
                        st.caption("（Excel(.xlsx)出力は openpyxl 未導入のため利用不可。CSVをご利用ください）")
                except Exception as e:
                    st.error(f"表の表示に失敗しました（CSV形式を確認してください）: {e}")
                    st.code(ws_data["code"], language="text")
                with st.expander("⚙️ CSV を編集"):
                    edited_code = st.text_area("CSV Editor", value=ws_data["code"], height=250)
                    if st.button("UPDATE SHEET"):
                        ws_data["code"] = edited_code
                        st.rerun()
            else:
                st.info("Canvas is empty.")

        elif ws_type == "DOC":
            if ws_data["code"]:
                st.success("✅ 文書が完成しました。Markdown / Word でダウンロードできます。")
                _safe = ws_name.replace(' ', '_')
                with st.container(border=True):
                    st.markdown(ws_data["code"])
                st.download_button("[ DOWNLOAD .md ]", data=ws_data["code"].encode("utf-8"),
                                   file_name=f"{_safe}.md", mime="text/markdown", use_container_width=True)
                try:
                    from docx import Document as _Docx
                    _doc = _Docx()
                    for _line in ws_data["code"].split("\n"):
                        _s = _line.rstrip()
                        if _s.startswith("### "): _doc.add_heading(_s[4:], level=3)
                        elif _s.startswith("## "): _doc.add_heading(_s[3:], level=2)
                        elif _s.startswith("# "): _doc.add_heading(_s[2:], level=1)
                        elif _s.startswith(("- ", "* ")): _doc.add_paragraph(_s[2:], style="List Bullet")
                        elif _s.strip(): _doc.add_paragraph(_s)
                    _dio = io.BytesIO(); _doc.save(_dio)
                    st.download_button("[ DOWNLOAD .docx ]", data=_dio.getvalue(),
                                       file_name=f"{_safe}.docx",
                                       mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                       type="primary", use_container_width=True)
                except Exception:
                    st.caption("（Word(.docx)出力は python-docx 未導入のため利用不可。Markdownをご利用ください）")
                with st.expander("⚙️ 本文(Markdown)を編集"):
                    edited_code = st.text_area("Markdown Editor", value=ws_data["code"], height=300)
                    if st.button("UPDATE DOC"):
                        ws_data["code"] = edited_code
                        st.rerun()
            else:
                st.info("Canvas is empty.")

    # ==========================================
    # AI実行ロジック（4つの脳みそ切り替え）
    # ==========================================
    trigger_prompt = forge_prompt if submitted else spoken_text if spoken_text else None
    
    is_auto_fix = False
    if st.session_state.auto_fix_prompt:
        trigger_prompt = st.session_state.auto_fix_prompt
        st.session_state.auto_fix_prompt = ""
        is_auto_fix = True

    if trigger_prompt:
        if is_auto_fix:
            ws_data["chat"].append({"role": "system", "avatar": "⚠️", "content": f"⚙️ AUTO-HEALING INITIATED:\n{trigger_prompt}"})
            sys_msg = "Auto-Healing in progress..."
        else:
            ws_data["chat"].append({"role": "user", "avatar": "👤", "content": trigger_prompt})
            ws_data["retries"] = 0
            sys_msg = f"Processing {ws_type.split(' ')[0]} Request..."

        with st.spinner(sys_msg):
            try:
                history_text = "【これまでの会話履歴】\n" + "\n".join([f"{msg['role']}: {msg['content']}" for msg in ws_data["chat"][:-1]])
                
                # 🧠 究極のシステムプロンプト（The Master Prompts）
                if "APP" in ws_type:
                    system_instruction = f"""
                    あなたは世界トップクラスのStreamlitアプリ開発者（シニアアーキテクト）です。
                    ユーザーの曖昧な指示から「本当に必要な機能」を推測し、見た目も美しく、バグのない完璧なPythonコードを提供します。
                    
                    【絶対遵守のレイアウト保護ルール】
                    1. `st.sidebar`、`st.set_page_config`、`st.chat_input` はOSの親画面を破壊するため【絶対に使用禁止】。チャットが必要な場合は `st.text_input` と `st.button` で代用せよ。
                    
                    【プロフェッショナルな開発要件】
                    1. 生成するアプリは完全独立型とし、`st.session_state` を活用して状態を適切に管理すること。
                    2. 単なる機能だけでなく、CSS（`st.markdown`）を駆使して「Neumorphism（ニューモーフィズム）」「Glassmorphism（グラスモーフィズム）」などのモダンで美しいUIを必ず実装すること。
                    3. 外部API（Gemini等）を使う場合は、必ずアプリ内に `st.text_input(..., type="password")` でキーを入力させる安全な設計にすること。
                    4. エラーハンドリング（`try-except`）を徹底し、ユーザーに優しいエラーメッセージ（`st.error` / `st.warning`）を表示すること。
                    5. コードは `# ...中略...` などの省略を絶対にせず、1行目から最後まで完全に出力すること。
                    
                    【出力フォーマット】
                    ```python
                    （ここに完全なコード）
                    ```
                    💡 次の拡張アイデア：
                    （アプリをより良くするためのプロ目線の提案を3つ）
                    
                    {history_text}
                    """
                
                elif "IMAGE" in ws_type:
                    system_instruction = f"""
                    あなたはMidjourneyやStable Diffusion、Imagen等の画像生成AIを完璧に操る、世界トップクラスの「AIプロンプトエンジニア兼アートディレクター」です。
                    ユーザーの簡単な日本語の要望から、画像生成AIが最も高品質で芸術的な画像を出力できる【究極の英語プロンプト】を構築します。
                    
                    【プロンプト構築の原則】
                    以下の要素をカンマ区切りの英語で緻密に記述すること：
                    1. Subject (主題): 構図、ポーズ、服装、表情
                    2. Medium (媒体): 写真、油絵、3Dレンダリング、水彩画、ベクターアートなど
                    3. Environment (環境): 背景、時間帯、天候、雰囲気
                    4. Lighting (照明): Cinematic lighting, volumetric lighting, rim lighting, soft softbox, neon lightsなど
                    5. Camera/Lens (カメラ設定): 35mm lens, f/1.8, macro photography, depth of field, drone shotなど
                    6. Style/Engine (スタイル): Unreal Engine 5, Octane Render, 8k resolution, highly detailed, masterpieceなど
                    
                    【出力フォーマット】
                    必ず以下の隠しタグ内に英語のプロンプトを記述すること。
                    [IMAGE_PROMPT: (ここに構築した緻密な英語プロンプト)]
                    
                    プロンプトの後に、日本語で「どのような意図でこのプロンプトを設計したか」の解説と、「さらに別のテイストにするためのアイデア」を簡潔に添えること。
                    
                    {history_text}
                    """
                    
                elif "VIDEO" in ws_type:
                    system_instruction = f"""
                    あなたはハリウッドで活躍する一流の映像ディレクター兼、SoraやVeoなどの最先端「動画生成AI」のプロンプトスペシャリストです。
                    ユーザーの要望から、プロの映像作品を作るための「絵コンテ構成」と、AIに直接入力する「英語の動画生成プロンプト」を作成します。
                    
                    【動画プロンプト（英語）の必須要素】
                    - Camera Movement (カメラワーク): Panning, Tilt, Tracking shot, Dolly zoom, FPV drone shot, Slow motionなど
                    - Scene Description (情景): 物理的な動き、光の反射、パーティクル（埃や火の粉）、被写界深度の変化
                    - Resolution/Style (画質): 4k, photorealistic, cinematic, 60fps
                    
                    【出力フォーマット（Markdown）】
                    ## 🎬 映像コンセプト
                    （どのような映像になるかの日本語解説）
                    
                    ## 🎥 シーン構成（絵コンテ）
                    - **Scene 1 (0:00-0:03):** （日本語でのシーン説明）
                    - **Scene 2 (0:03-0:06):** （日本語でのシーン説明）
                    ...
                    
                    ## 🤖 AI用動画生成プロンプト (English)
                    `（ここに動画生成AIにそのままコピペできる、すべてのシーンを統合した高品質な英語プロンプトを記述）`
                    
                    {history_text}
                    """
                    
                elif "SLIDE" in ws_type:
                    system_instruction = f"""
                    あなたはマッキンゼーなどのトップコンサルティングファームで活躍する、一流のプレゼン・ストラテジストです。
                    ユーザーのテーマに基づき、聴衆を惹きつける論理的で説得力のあるプレゼンテーション資料（スライド構成）を作成します。
                    
                    【スライド構築の絶対ルール】
                    1. 構成は「結序破急」または「PREP法」など、論理的なストーリーテリングを意識すること。
                    2. スライドとスライドの間は必ず `---` (ハイフン3つ) のみで区切ること（システムがこれでスライドを分割します）。
                    3. 1枚のスライドの情報量は多すぎず、視覚的に分かりやすいMarkdown（箇条書き、太字、引用）を使うこと。
                    
                    【出力フォーマット例】
                    # スライドタイトル1
                    ### サブタイトル
                    - ポイントA
                    - ポイントB
                    > 印象的な引用やメッセージ
                    
                    ---
                    
                    # スライドタイトル2
                    ...
                    
                    スライドを出力した後に、「このプレゼンを話す際のアドバイス（トークスクリプトのヒント）」を日本語で添えること。

                    {history_text}
                    """

                elif "SHEET" in ws_type:
                    system_instruction = f"""
                    あなたは一流のデータアナリスト兼スプレッドシート設計者です。
                    ユーザーの要望から、実用的な表データ（スプレッドシート）を設計して出力します。

                    【出力ルール】
                    1. 表データは必ずCSV形式とし、先頭行をヘッダーにして、次のコードフェンス内にのみ記述すること：
                    ```csv
                    列1,列2,列3
                    値,値,値
                    ```
                    2. 数値は単位や桁区切りを入れず、計算可能な生の数値で出力すること。
                    3. 行数は要望に応じて十分な実用量（目安5〜30行）を用意すること。

                    CSVの後に、日本語で「この表の使い方」と「追加すると便利な列」を簡潔に添えること。

                    {history_text}
                    """

                elif "DOC" in ws_type:
                    system_instruction = f"""
                    あなたはプロのビジネスライター兼編集者です。
                    ユーザーの要望から、構造化された読みやすい文書をMarkdownで作成します。

                    【出力ルール】
                    1. 見出し(#, ##, ###)、箇条書き(-)、太字(**)、表、引用(>)を適切に使い、論理的な構成にすること。
                    2. そのまま提出できる完成度で、省略（…中略…）をしないこと。
                    3. 文書本体のみを出力し、余計な前置きは書かないこと。

                    {history_text}
                    """

                # 🤖 マルチAI対応：get_ai_response 経由で呼び出す（Gemini→Claude→Grok→OpenAI）
                ai_text = get_ai_response(system_instruction + "\n\nボスの現在の指示: " + trigger_prompt, model='gemini-2.5-flash')
                
                # モードごとの後処理（パース）
                reply_text = ai_text
                
                if "APP" in ws_type:
                    code_match = re.search(r'```python\n(.*?)\n```', ai_text, re.DOTALL)
                    if code_match:
                        ws_data["code"] = code_match.group(1)
                        reply_text = ai_text.replace(code_match.group(0), "").strip() or "アプリケーションのコードを構築しました。"
                
                elif "IMAGE" in ws_type:
                    prompt_match = re.search(r'\[IMAGE_PROMPT:\s*(.*?)\]', ai_text)
                    if prompt_match:
                        ws_data["media"] = prompt_match.group(1).strip()
                        reply_text = ai_text.replace(prompt_match.group(0), "").strip()
                
                elif "SHEET" in ws_type:
                    csv_match = re.search(r'```csv\n(.*?)\n```', ai_text, re.DOTALL)
                    if csv_match:
                        ws_data["code"] = csv_match.group(1).strip()
                        reply_text = ai_text.replace(csv_match.group(0), "").strip() or "表データを作成しました。"
                    else:
                        ws_data["code"] = ai_text  # フェンスが無い場合もそのまま保持
                        reply_text = "表データを作成しました。"

                elif "DOC" in ws_type:
                    ws_data["code"] = ai_text
                    reply_text = "文書を作成しました。右のプレビューを確認してください。"

                elif "VIDEO" in ws_type or "SLIDE" in ws_type:
                    ws_data["code"] = ai_text
                    reply_text = "資料の作成が完了しました。右のプレビュー画面を確認してください。"

                # 音声の生成と保存
                tts = gTTS(text=re.sub(r'[*#`_]', '', reply_text[:200]), lang='ja')
                audio_fp = io.BytesIO()
                tts.write_to_fp(audio_fp)
                st.session_state.ai_voice_base64 = base64.b64encode(audio_fp.getvalue()).decode()
                st.session_state.just_generated_audio = True
                
                ws_data["chat"].append({"role": "assistant", "avatar": "🤖", "content": reply_text})
                st.rerun()
            except Exception as e:
                st.error(f"Error: {e}")

# ==========================================
# 📥 アプリ保存モジュール (APPモード時のみ表示)
# ==========================================
if st.session_state.current_forge_ws and "APP" in ws_data.get("type", ""):
    st.markdown("---")
    st.markdown("#### 💾 SAVE TO APP ARCHIVE")
    with st.expander("📦 新しいミニアプリとしてインストール", expanded=False):
        app_filename = st.text_input("アプリのファイル名（半角英数字）", placeholder="例: my_calculator")
        app_code_input = st.text_area("保存するPythonコードを貼り付け", height=250, value=ws_data.get("code", ""))
        
        if st.button("ARCHIVE にインストール ⚡", use_container_width=True, type="primary"):
            if app_filename and app_code_input:
                safe_name = app_filename.replace(" ", "_").lower()
                if not safe_name.endswith(".py"): safe_name += ".py"
                os.makedirs("forge_apps", exist_ok=True)
                try:
                    with open(os.path.join("forge_apps", safe_name), "w", encoding="utf-8") as f:
                        f.write(app_code_input)
                    _apps = vault_get("forge_apps", {}) or {}
                    _apps[safe_name] = app_code_input
                    persist_vault_key("forge_apps", _apps)  # 再起動でも残るようVaultにも保存
                    st.success(f"✅ インストール完了！ `{safe_name}` をAPP ARCHIVEに保存しました。")
                    st.balloons()
                except Exception as e:
                    st.error(f"保存エラー: {e}")
